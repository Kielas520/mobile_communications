#!/usr/bin/env python3
"""移动通信原理复习资料提取工具。

用法:
    python3 scripts/extract_materials.py          # 提取全部 (pptx + pdf + heic)
    python3 scripts/extract_materials.py pptx      # 只提取 source/课程/*.pptx 文字
    python3 scripts/extract_materials.py pdf       # 只提取 source/复习资料/*.pdf 文字
    python3 scripts/extract_materials.py heic      # 只把 source/试卷/*.heic 转成 PNG

依赖:
    python-pptx (PPTX)、PyMuPDF/fitz (PDF)、macOS 自带 sips (HEIC 转 PNG)
"""
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SOURCE = ROOT / "source"
COURSE_DIR = SOURCE / "课程"
REVIEW_DIR = SOURCE / "复习资料"
EXAM_DIR = SOURCE / "试卷"
EXTRACTED = SOURCE / "extracted"


def extract_pptx() -> None:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("缺少依赖: 请先运行 `pip install python-pptx`")

    files = sorted(COURSE_DIR.glob("*.pptx"))
    if not files:
        print(f"未在 {COURSE_DIR} 找到 .pptx 文件")
        return
    for pptx_path in files:
        prs = Presentation(pptx_path)
        lines = [f"# {pptx_path.stem}\n"]
        for idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"\n## 幻灯片 {idx}\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        lines.append(" | ".join(cells))
        out_dir = EXTRACTED / pptx_path.stem
        out_dir.mkdir(parents=True, exist_ok=True)
        (out_dir / "content.md").write_text("\n".join(lines), encoding="utf-8")
        print(f"[PPTX] {pptx_path.name} -> {out_dir / 'content.md'}")


def extract_pdf() -> None:
    try:
        import fitz
    except ImportError:
        sys.exit("缺少依赖: 请先运行 `pip install PyMuPDF`")

    files = sorted(REVIEW_DIR.glob("*.pdf"))
    if not files:
        print(f"未在 {REVIEW_DIR} 找到 .pdf 文件")
        return
    for pdf_path in files:
        if pdf_path.stat().st_size == 0:
            print(f"[PDF] 跳过空文件: {pdf_path.name}")
            continue
        try:
            doc = fitz.open(pdf_path)
        except Exception as exc:
            print(f"[PDF] 无法打开 {pdf_path.name}: {exc}")
            continue
        page_texts = [page.get_text().strip() for page in doc]
        total_chars = sum(len(t) for t in page_texts)
        out_dir = EXTRACTED / pdf_path.stem
        out_dir.mkdir(parents=True, exist_ok=True)

        if total_chars >= 50:
            lines = [f"# {pdf_path.stem}\n"]
            for page_num, text in enumerate(page_texts, start=1):
                lines.append(f"\n## 第 {page_num} 页\n")
                lines.append(text)
            (out_dir / "content.md").write_text("\n".join(lines), encoding="utf-8")
            print(f"[PDF] {pdf_path.name} -> {out_dir / 'content.md'} (文字版)")
        else:
            img_dir = out_dir / "页面图片"
            img_dir.mkdir(parents=True, exist_ok=True)
            lines = [
                f"# {pdf_path.stem}\n",
                "> 本 PDF 为图片型（扫描/图片幻灯片），无可提取文字。已逐页渲染为 PNG，请用 `read` 工具直接看图。\n",
            ]
            mat = fitz.Matrix(2, 2)
            for page_num, page in enumerate(doc, start=1):
                png_path = img_dir / f"P{page_num:02d}.png"
                page.get_pixmap(matrix=mat).save(png_path)
                lines.append(f"\n## 第 {page_num} 页\n")
                lines.append(f"![第{page_num}页](页面图片/P{page_num:02d}.png)")
            (out_dir / "content.md").write_text("\n".join(lines), encoding="utf-8")
            print(
                f"[PDF] {pdf_path.name} 为图片型 -> 渲染 {doc.page_count} 页到 {img_dir}"
            )
        doc.close()


def convert_heic() -> None:
    files = sorted(EXAM_DIR.glob("*.heic")) + sorted(EXAM_DIR.glob("*.HEIC"))
    if not files:
        print(f"未在 {EXAM_DIR} 找到 .heic 文件")
        return
    out_dir = EXAM_DIR / "png"
    out_dir.mkdir(parents=True, exist_ok=True)
    for heic_path in files:
        out_path = out_dir / (heic_path.stem + ".png")
        subprocess.run(
            ["sips", "-s", "format", "png", str(heic_path), "--out", str(out_path)],
            check=True,
            stdout=subprocess.DEVNULL,
        )
        print(f"[HEIC] {heic_path.name} -> {out_path}")


def _rel(p: Path) -> str:
    return p.relative_to(ROOT).as_posix()


def scaffold() -> None:
    """为图片类资料生成"待转录"目标文件与任务清单。

    实际转录需由支持视觉的模型读取图片后填入；本步骤只搭骨架。
    """
    categories = [
        (
            "马志乐手写复习笔记（最高优先级，先读这个）",
            sorted((EXTRACTED / "移动通信复习-马志乐" / "页面图片").glob("P*.png")),
            EXTRACTED / "移动通信复习-马志乐" / "手写笔记转录.md",
        ),
        (
            "平时作业答案",
            sorted((SOURCE / "作业答案").glob("*.png")),
            EXTRACTED / "作业答案" / "转录.md",
        ),
        (
            "各章小结习题答案",
            sorted((SOURCE / "移动通信小结答案").glob("*.png")),
            EXTRACTED / "移动通信小结答案" / "转录.md",
        ),
        (
            "往期试卷",
            sorted((EXAM_DIR / "png").glob("*.png")),
            EXTRACTED / "试卷" / "转录.md",
        ),
    ]

    checklist = ["# 图片转录任务清单\n",
                 "> 当前（无视觉）模型无法读图。换成支持视觉的模型后，逐张读取下列图片，"
                 "把题目/答案/公式/图示说明转录进各自的目标文件。**马志乐手写笔记最先读。**\n"]

    for label, images, target in categories:
        target.parent.mkdir(parents=True, exist_ok=True)
        if not images:
            print(f"[SCAFFOLD] 跳过（无图片）: {label}")
            continue
        if target.exists():
            print(f"[SCAFFOLD] 已存在，保留不覆盖: {_rel(target)}")
        else:
            body = [f"# {label} — 转录\n",
                    "> 读图后把每张图的内容填到对应小节，删掉「待转录」占位。\n"]
            for img in images:
                body.append(f"\n## {img.name}\n")
                body.append(f"图片路径：`{_rel(img)}`\n")
                body.append("待转录（题目 / 答案 / 公式 / 图示说明）\n")
            target.write_text("\n".join(body), encoding="utf-8")
            print(f"[SCAFFOLD] {label}: {len(images)} 张 -> {_rel(target)}")

        checklist.append(f"\n## {label}（{len(images)} 张）")
        checklist.append(f"- 目标文件：`{_rel(target)}`")
        for img in images:
            checklist.append(f"  - [ ] `{_rel(img)}`")

    (EXTRACTED / "图片转录任务清单.md").write_text("\n".join(checklist), encoding="utf-8")
    print(f"[SCAFFOLD] 任务清单 -> {_rel(EXTRACTED / '图片转录任务清单.md')}")


def main() -> None:
    mode = sys.argv[1] if len(sys.argv) > 1 else "all"
    if mode in ("all", "pptx"):
        extract_pptx()
    if mode in ("all", "pdf"):
        extract_pdf()
    if mode in ("all", "heic"):
        convert_heic()
    if mode in ("all", "scaffold"):
        scaffold()
    if mode not in ("all", "pptx", "pdf", "heic", "scaffold"):
        sys.exit(f"未知模式: {mode}（可选: all/pptx/pdf/heic/scaffold）")


if __name__ == "__main__":
    main()
